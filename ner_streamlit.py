"""
Organization Extraction Streamlit Application
=============================================

Interactive web application for analysts to upload documents,
extract organizations, and manage database additions.
"""

import streamlit as st
import pandas as pd
import numpy as np
import re
import time
import json
import os
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from datetime import datetime
import psycopg2
from sqlalchemy import create_engine, text
from transformers import AutoTokenizer, pipeline
from fuzzywuzzy import fuzz, process
from collections import defaultdict, Counter
import warnings
warnings.filterwarnings('ignore')
import streamlit.components.v1 as components

# Fix torch compatibility issue with Streamlit
import os
os.environ["TOKENIZERS_PARALLELISM"] = "false"

# Load environment variables
from dotenv import load_dotenv
load_dotenv()

# Import document processing libraries
import PyPDF2
import docx
from io import BytesIO
import openpyxl
from pptx import Presentation

class OrganizationExtractor:
    """Core organization extraction system for Streamlit app"""
    
    def __init__(self, postgres_config: Dict[str, str]):
        self.postgres_config = postgres_config
        self.min_confidence = 0.85  # Conservative threshold for production
        self.min_org_length = 3
        
        # Initialize components
        self.master_orgs_df = pd.DataFrame()
        self.org_lookup = {}
        self.tokenizer = None
        self.ner_model = None
        
        # Initialize system
        self._initialize_system()
    
    @st.cache_resource
    def _initialize_system(_self):
        """Initialize system components with caching"""
        try:
            # Initialize tokenizer
            _self.tokenizer = AutoTokenizer.from_pretrained("bert-base-uncased")
            
            # Load NER model
            _self.ner_model = pipeline(
                "ner",
                model="dbmdz/bert-large-cased-finetuned-conll03-english",
                aggregation_strategy="simple",
                device=-1
            )
            
            # Connect to database and load organizations
            _self._load_organizations()
            _self._build_lookup()
            
            return True
            
        except Exception as e:
            st.error(f"System initialization failed: {e}")
            return False
    
    def _load_organizations(self):
        """Load organizations from database"""
        try:
            conn = psycopg2.connect(
                host=self.postgres_config['host'],
                port=self.postgres_config.get('port', 5432),
                database=self.postgres_config['database'],
                user=self.postgres_config['user'],
                password=self.postgres_config['password']
            )
            
            query = """
            SELECT org_id, org_name
            FROM verdantix.org
            WHERE org_name IS NOT NULL 
                AND LENGTH(TRIM(org_name)) > 2
            ORDER BY org_name
            """
            
            self.master_orgs_df = pd.read_sql_query(query, conn)
            conn.close()
            
        except Exception as e:
            st.error(f"Database connection failed: {e}")
            self.master_orgs_df = pd.DataFrame()
    
    def _build_lookup(self):
        """Build organization lookup dictionary"""
        self.org_lookup = {}
        self.known_short_orgs = set()  # Build dynamically from database
        
        for _, org in self.master_orgs_df.iterrows():
            org_id = org['org_id']
            org_name = str(org['org_name']).strip()
            
            if len(org_name) > 2:
                self.org_lookup[org_name.lower()] = {
                    'org_id': org_id,
                    'canonical': org_name,
                    'confidence': 1.0
                }
                
                # Track short orgs from database
                if len(org_name) < 8 and ' ' not in org_name:
                    self.known_short_orgs.add(org_name.lower())
                
                # Add aliases
                aliases = self._generate_aliases(org_name)
                for alias in aliases:
                    if alias.lower() not in self.org_lookup and len(alias) > 2:
                        self.org_lookup[alias.lower()] = {
                            'org_id': org_id,
                            'canonical': org_name,
                            'confidence': 0.85
                        }
                        
                        # Track short aliases too
                        if len(alias) < 8 and ' ' not in alias:
                            self.known_short_orgs.add(alias.lower())
    
    def _generate_aliases(self, org_name: str) -> List[str]:
        """Generate aliases for organization names"""
        aliases = []
        
        # Extract acronyms from parentheses (e.g., "World Economic Forum (WEF)" -> "WEF")
        paren_pattern = r'\(([^)]+)\)'
        paren_matches = re.findall(paren_pattern, org_name)
        for match in paren_matches:
            cleaned = match.strip()
            if 2 <= len(cleaned) <= 10:
                aliases.append(cleaned)
        
        # Remove business suffixes
        suffix_patterns = [
            r'\s+(?:Inc\.?|Corporation|Corp\.?|Company|Co\.?|Limited|Ltd\.?)',
            r'\s+(?:LLC|LLP|LP|PLC|Group|Holdings?)'
        ]
        
        for pattern in suffix_patterns:
            base = re.sub(pattern + r'$', '', org_name, flags=re.IGNORECASE).strip()
            if base != org_name and len(base) > 2:
                aliases.append(base)
        
        # Create acronyms from main text (excluding parentheses)
        main_text = re.sub(paren_pattern, '', org_name).strip()
        words = main_text.split()
        if len(words) > 1:
            # Skip common words when creating acronyms
            skip_words = {'of', 'the', 'and', 'for', 'in', 'on', 'at', 'to', 'a', 'an', '&'}
            meaningful_words = [w for w in words if w.lower() not in skip_words and len(w) > 0]
            if len(meaningful_words) > 1:
                acronym = ''.join([w[0].upper() for w in meaningful_words])
                if 2 <= len(acronym) <= 8:
                    aliases.append(acronym)
        
        return list(set(aliases))
    
    def _is_generic_term(self, term: str) -> bool:
        """Filter generic business terms"""
        generic_terms = {
            'ai', 'iot', 'esg', 'api', 'cloud', 'data', 'tech', 'digital',
            'smart', 'green', 'cyber', 'auto', 'bio', 'blockchain', 'fintech',
            'saas', 'crm', 'erp', 'hr', 'it', 'covid', 'gdpr',
            # Add common words that shouldn't be organizations
            'are', 'is', 'was', 'were', 'be', 'been', 'have', 'has', 'had',
            'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might',
            'can', 'must', 'shall', 'to', 'of', 'in', 'for', 'on', 'at', 'by',
            'with', 'from', 'up', 'about', 'into', 'through', 'during', 'before',
            'after', 'above', 'below', 'under', 'over', 'between', 'among',
            'security', 'chain', 'track', 'progress', 'time', 'area', 'areas',
            'effect', 'impact', 'data', 'information', 'report', 'reports',
            'standard', 'standards', 'framework', 'frameworks'
        }
        return term.lower() in generic_terms

    def _is_false_positive_prone(self, term: str) -> bool:
        """Check if term is prone to false positives and needs stricter validation"""
        # Terms that appear in many contexts but could be org names
        false_positive_prone = {
            'sustainability', 'sustainibility', 'innovation', 'development',
            'technology', 'solutions', 'services', 'management', 'consulting',
            'operations', 'strategic', 'digital', 'transformation', 'analytics',
            'research', 'institute', 'group', 'partners', 'systems', 'network', 
            'accountability', 'benchmark', 'framework', 'automation', 'measurable',
            'decisions', 'interface', 'frontier','overview', 'seamless','succeed',
            'access', 'deploy', 'engage', 'fabric', 'figure', 'switch', 'aware', 
            'focus', 'given', 'pilot', 'trend', 'best', 'code', 'edge', 'here', 'next', 
            'true', 'well', 'box', 'vector', 'group', 'space', 'partners', 'consultants', 
            'ecosystem', 'reliance', 'current', 'nature', 'boots', 'first', 'reach', 'shape', 
            'near', 'peak', 'view', 'era', 'gap', 'one', 'Australia', 'capital', 'London', 'core', 
            'five', 'post', 'slam', 'big', 'cog', 'its', 'net', 'PLC', 'RFP'
        }
        return any(prone_term in term.lower() for prone_term in false_positive_prone)

    def _requires_strong_context(self, term: str) -> bool:
        """Check if term requires strong organizational context due to ambiguity"""
        # Single words that are highly ambiguous
        if ' ' not in term:
            highly_ambiguous = {
                'sustainability', 'innovation', 'development', 'technology',
                'management', 'consulting', 'operations', 'research', 'analytics', 
                'accountability','benchmark', 'framework', 'automation', 'measurable',
                'decisions', 'interface', 'frontier', 'overview', 'seamless','succeed',
                'access','deploy', 'engage', 'fabric', 'figure', 'switch', 'aware', 
                'focus', 'given', 'pilot', 'trend', 'best', 'code', 'edge', 'here', 'next', 
                'true', 'well', 'box', 'vector', 'group', 'space', 'partners', 'consultants', 
                'ecosystem', 'reliance', 'current', 'nature', 'boots', 'first', 'reach', 'shape', 
                'near', 'peak', 'view', 'era', 'gap', 'one', 'Australia', 'capital', 'London', 'core', 
                'five', 'post', 'slam', 'big', 'cog', 'its', 'net', 'PLC', 'RFP'
            }
            return term.lower() in highly_ambiguous

        # Multi-word phrases containing ambiguous terms
        return self._is_false_positive_prone(term)

    def _get_confidence_penalty(self, term: str) -> float:
        """Apply confidence penalty for potentially ambiguous terms"""
        if self._is_false_positive_prone(term):
            # Higher penalty for single words that are very common
            if ' ' not in term and term.lower() in ['sustainability', 'innovation', 'development', 'accountability','benchmark', 'framework', 'automation', 'measurable',
                                                    'decisions', 'interface', 'frontier', 'overview', 'seamless','succeed','access','deploy', 'engage', 'fabric',
                                                    'figure', 'switch', 'aware', 
                                                    'focus', 'given', 'pilot', 'trend', 'best', 'code', 'edge', 'here', 'next', 
                                                     'true', 'well', 'box', 'vector', 'group', 'space', 'partners', 'consultants', 'ecosystem', 'reliance', 'current',
                                                    'nature', 'boots', 'first', 'reach', 'shape', 'near', 'peak', 'view', 'era', 'gap', 'one', 'Australia', 'capital',
                                                      'London', 'core', 'five', 'post', 'slam', 'big', 'cog', 'its', 'net', 'PLC', 'RFP']:
                return 0.4  # 40% penalty
            else:
                return 0.15  # 15% penalty
        return 0.0
    
    def _validate_context(self, text: str, start: int, end: int, term: str = None) -> bool:
        """Validate organizational context with enhanced filtering for ambiguous terms"""

        # For ambiguous terms, require stronger context
        if term and self._requires_strong_context(term):
            return self._validate_strong_context(text, start, end)

        # Standard context validation
        window = 100
        context_start = max(0, start - window)
        context_end = min(len(text), end + window)
        context = text[context_start:context_end].lower()

        org_indicators = [
            'company', 'corporation', 'inc', 'founded', 'ceo', 'announced',
            'partnership', 'acquisition', 'investment', 'subsidiary', 'firm', 'firms',
            'business', 'enterprise', 'organization', 'organisation', 'group', 'ltd',
            'limited', 'consulting', 'consultancy', 'services', 'solutions', 'technologies',
            'technology', 'corporate', 'management', 'strategic', 'operations', 'client',
            'clients', 'market', 'industry', 'sector', 'expertise', 'capabilities',
            'team', 'teams', 'office', 'offices', 'division', 'department', 'unit',
            'board', 'executive', 'director', 'manager', 'leadership', 'staff',
            'employee', 'employees', 'workforce', 'personnel', 'professional', 'specialist'
        ]

        return any(indicator in context for indicator in org_indicators)

    def _validate_strong_context(self, text: str, start: int, end: int) -> bool:
        """Stricter context validation for ambiguous terms"""
        window = 50  # Smaller window for stricter validation
        context_start = max(0, start - window)
        context_end = min(len(text), end + window)
        context = text[context_start:context_end].lower()

        # Require very explicit organizational indicators
        strong_org_indicators = [
            'company', 'corporation', 'inc', 'inc.', 'ltd', 'ltd.', 'llc',
            'founded', 'ceo', 'headquarters', 'subsidiary', 'acquired',
            'partnership with', 'announced by', 'reported by', 'according to',
            'spokesperson for', 'representative from', 'director of'
        ]

        return any(indicator in context for indicator in strong_org_indicators)

    def _find_organization_patterns(self, text: str) -> List[Dict]:
        """Find organization patterns like 'Full Name (ACRONYM)' that NER might miss"""
        patterns = []

        # Pattern: "Full Organization Name (ACRONYM)" - improved to avoid greedy matching
        pattern = r'(?:^|(?<=\s))([A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z]+)*)\s*\(([A-Z]{2,6})\)'

        for match in re.finditer(pattern, text):
            full_name = match.group(1).strip()
            acronym = match.group(2).strip()
            start, end = match.span()

            # Check if this looks like an organization
            if (len(full_name.split()) >= 2 and  # Multi-word
                not any(word.lower() in ['the', 'of', 'and', 'for', 'in'] for word in full_name.split()[0:1]) and  # Doesn't start with article
                len(acronym) >= 2):  # Valid acronym

                patterns.append({
                    'full_text': match.group(0),
                    'full_name': full_name,
                    'acronym': acronym,
                    'start': start,
                    'end': end,
                    'confidence': 0.9
                })

        return patterns

    def extract_organizations(self, text: str) -> Tuple[List[Dict], List[Dict]]:
        """Extract organizations and categorize as known vs unknown"""
        
        db_matches = []
        ner_discoveries = []
        
        # Method 1: Database matches
        text_lower = text.lower()
        sorted_terms = sorted(self.org_lookup.keys(), key=len, reverse=True)
        matched_positions = set()

        # Pre-process text to find organization patterns like "Name (ACRONYM)"
        org_patterns = self._find_organization_patterns(text)

        # Handle detected patterns first
        for pattern in org_patterns:
            # Check if the FULL PATTERN TEXT is in database first (highest priority)
            full_text_lower = pattern['full_text'].lower()
            full_name_lower = pattern['full_name'].lower()
            acronym_lower = pattern['acronym'].lower()

            if full_text_lower in self.org_lookup:
                # Exact full text match (e.g., "Global Reporting Initiative (GRI)")
                org_info = self.org_lookup[full_text_lower]
                db_matches.append({
                    'text': pattern['full_text'],
                    'canonical': org_info['canonical'],
                    'confidence': org_info['confidence'],
                    'org_id': org_info['org_id'],
                    'method': 'database'  # This is a direct database match
                })
                matched_positions.add((pattern['start'], pattern['end']))

            elif full_name_lower in self.org_lookup:
                # Full name match (without parentheses)
                org_info = self.org_lookup[full_name_lower]
                db_matches.append({
                    'text': pattern['full_text'],
                    'canonical': org_info['canonical'],
                    'confidence': org_info['confidence'],
                    'org_id': org_info['org_id'],
                    'method': 'pattern_full'
                })
                matched_positions.add((pattern['start'], pattern['end']))

            elif acronym_lower in self.org_lookup:
                # Acronym match (lowest priority)
                org_info = self.org_lookup[acronym_lower]
                db_matches.append({
                    'text': pattern['full_text'],
                    'canonical': f"{pattern['full_name']} ({pattern['acronym']})",
                    'confidence': pattern['confidence'],
                    'org_id': org_info['org_id'],
                    'method': 'pattern_acronym'
                })
                matched_positions.add((pattern['start'], pattern['end']))

            else:
                # New organization pattern
                ner_discoveries.append({
                    'text': pattern['full_text'],
                    'confidence': pattern['confidence'],
                    'method': 'pattern_new',
                    'full_name': pattern['full_name'],
                    'acronym': pattern['acronym']
                })

        for term in sorted_terms:
            if (len(term) >= self.min_org_length and
                not self._is_generic_term(term)):

                # Require multi-word for short terms, but allow known brands from database
                if len(term) < 8 and ' ' not in term:
                    # Use dynamically built list of short orgs from database
                    if term.lower() not in self.known_short_orgs:
                        continue

                # Handle parentheses in organization names - word boundaries don't work well with them
                if '(' in term or ')' in term:
                    # Use lookahead/lookbehind for terms with parentheses
                    pattern = r'(?<!\w)' + re.escape(term) + r'(?!\w)'
                else:
                    # Standard word boundary for terms without parentheses
                    pattern = r'\b' + re.escape(term) + r'\b'

                for match in re.finditer(pattern, text_lower):
                    start, end = match.span()

                    if not any(start < e and s < end for s, e in matched_positions):
                        matched_positions.add((start, end))

                        org_info = self.org_lookup[term]

                        # Apply context validation for false-positive prone terms
                        if self._is_false_positive_prone(term):
                            if not self._validate_context(text, start, end, term):
                                continue  # Skip this match if context doesn't support it

                        # Apply confidence penalty for ambiguous terms
                        base_confidence = org_info['confidence']
                        penalty = self._get_confidence_penalty(term)
                        final_confidence = max(0.1, base_confidence - penalty)

                        # Only add if confidence is above minimum threshold
                        if final_confidence >= 0.7:  # Require reasonable confidence for DB matches
                            db_matches.append({
                                'text': text[start:end],
                                'canonical': org_info['canonical'],
                                'confidence': final_confidence,
                                'org_id': org_info['org_id'],
                                'method': 'database',
                                'had_penalty': penalty > 0
                            })
        
        # Method 2: NER for new organizations
        if self.ner_model:
            try:
                if len(text) > 4000:
                    chunks = [text[i:i+4000] for i in range(0, len(text), 3500)]
                else:
                    chunks = [text]
                
                for chunk in chunks:
                    predictions = self.ner_model(chunk)
                    
                    for pred in predictions:
                        if 'ORG' in pred.get('entity_group', ''):
                            org_text = pred['word'].strip()
                            
                            # Clean tokenization artifacts
                            org_text = re.sub(r'^##', '', org_text)
                            org_text = re.sub(r'[^\w\s&.-]', '', org_text)
                            org_text = ' '.join(org_text.split())
                            
                            if (len(org_text) >= self.min_org_length and
                                pred['score'] >= self.min_confidence and
                                not self._is_generic_term(org_text)):

                                # Apply confidence penalty for ambiguous terms
                                base_confidence = pred['score']
                                penalty = self._get_confidence_penalty(org_text)
                                final_confidence = max(0.1, base_confidence - penalty)

                                # Skip if confidence drops too low after penalty
                                if final_confidence < self.min_confidence:
                                    continue

                                # Check if already in database
                                if org_text.lower() not in self.org_lookup:
                                    # Try fuzzy matching with consistent results
                                    if len(self.master_orgs_df) > 0:
                                        canonical_names = self.master_orgs_df['org_name'].tolist()
                                        
                                        # Sort canonical names for deterministic results
                                        canonical_names = sorted(canonical_names)
                                        
                                        best_match = process.extractOne(
                                            org_text, canonical_names, scorer=fuzz.ratio
                                        )
                                        
                                        # Store fuzzy match info for debugging
                                        fuzzy_score = best_match[1] if best_match else 0
                                        
                                        if best_match and best_match[1] >= 85:
                                            # Close match found - add to DB matches
                                            matched_row = self.master_orgs_df[
                                                self.master_orgs_df['org_name'] == best_match[0]
                                            ].iloc[0]
                                            
                                            # Apply penalty to fuzzy match confidence too
                                            fuzzy_confidence = final_confidence * (best_match[1] / 100)

                                            db_matches.append({
                                                'text': org_text,
                                                'canonical': matched_row['org_name'],
                                                'confidence': fuzzy_confidence,
                                                'org_id': matched_row['org_id'],
                                                'method': 'ner_fuzzy',
                                                'fuzzy_score': fuzzy_score,
                                                'had_penalty': penalty > 0
                                            })
                                        else:
                                            # New organization discovery - include fuzzy match info
                                            ner_discoveries.append({
                                                'text': org_text,
                                                'confidence': final_confidence,
                                                'method': 'ner_new',
                                                'best_fuzzy_match': best_match[0] if best_match else None,
                                                'fuzzy_score': fuzzy_score,
                                                'had_penalty': penalty > 0,
                                                'penalty_amount': penalty
                                            })
            
            except Exception as e:
                st.error(f"NER processing error: {e}")
        
        # Deduplicate
        db_matches = self._deduplicate_matches(db_matches)
        ner_discoveries = self._deduplicate_matches(ner_discoveries)
        
        return db_matches, ner_discoveries
    
    def _deduplicate_matches(self, matches: List[Dict]) -> List[Dict]:
        """Remove duplicate matches"""
        seen = set()
        deduplicated = []
        
        for match in sorted(matches, key=lambda x: x['confidence'], reverse=True):
            canonical = match.get('canonical', match['text']).lower()
            if canonical not in seen:
                seen.add(canonical)
                deduplicated.append(match)
        
        return deduplicated


def extract_text_from_file(uploaded_file) -> str:
    """Extract text from uploaded file"""
    file_type = uploaded_file.type
    
    try:
        if file_type == "text/plain":
            return str(uploaded_file.read(), "utf-8")
        
        elif file_type == "application/pdf":
            pdf_reader = PyPDF2.PdfReader(BytesIO(uploaded_file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document(BytesIO(uploaded_file.read()))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        
        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # Excel files (.xlsx)
            workbook = openpyxl.load_workbook(BytesIO(uploaded_file.read()), data_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    if row_text:
                        text += " | ".join(row_text) + "\n"
                text += "\n"
            
            return text
        
        elif file_type == "application/vnd.ms-excel":
            # Legacy Excel files (.xls) - handled by openpyxl as well
            try:
                workbook = openpyxl.load_workbook(BytesIO(uploaded_file.read()), data_only=True)
                text = ""
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text += f"Sheet: {sheet_name}\n"
                    
                    for row in sheet.iter_rows():
                        row_text = []
                        for cell in row:
                            if cell.value is not None:
                                row_text.append(str(cell.value))
                        if row_text:
                            text += " | ".join(row_text) + "\n"
                    text += "\n"
                
                return text
            except Exception as e:
                st.error(f"Error reading legacy Excel file: {e}")
                return ""
        
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # PowerPoint files (.pptx)
            prs = Presentation(BytesIO(uploaded_file.read()))
            text = ""
            
            for i, slide in enumerate(prs.slides, 1):
                text += f"Slide {i}:\n"
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                
                # Extract text from slide notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        text += f"Notes: {notes_text}\n"
                
                text += "\n"
            
            return text
        
        else:
            st.error(f"Unsupported file type: {file_type}")
            return ""
            
    except Exception as e:
        st.error(f"Error reading file: {e}")
        return ""


def add_keep_alive_functionality():
    """Add JavaScript to prevent app from sleeping"""
    keep_alive_js = """
    <script>
        // Auto-refresh every 25 minutes (1500 seconds) to prevent sleeping
        setInterval(function() {
            // Create a small ping to keep the session alive
            fetch(window.location.href, {
                method: 'HEAD',
                cache: 'no-cache'
            }).catch(function(error) {
                console.log('Keep-alive ping failed:', error);
            });
        }, 1500000); // 25 minutes

        // Periodic activity simulation every 5 minutes
        setInterval(function() {
            // Trigger a small state change to maintain session
            const event = new Event('streamlit:stateChanged');
            window.dispatchEvent(event);
        }, 300000); // 5 minutes
    </script>
    """
    components.html(keep_alive_js, height=0)

def main():
    """Main Streamlit application"""

    st.set_page_config(
        page_title="Organization Extraction Tool",
        page_icon="🏢",
        layout="wide"
    )

    # Add keep-alive functionality
    add_keep_alive_functionality()

    st.title("Organization Extraction & Database Management")
    st.markdown("Upload documents to automatically extract and categorize organizations")
    
    # Configuration from Streamlit secrets (fallback to environment variables)
    try:
        # Check if secrets are available
        if hasattr(st.secrets, '_secrets') and "database" in st.secrets:
            # Use Streamlit secrets
            db_secrets = st.secrets["database"]
            postgres_config = {
                'host': db_secrets.get("POSTGRES_HOST", os.getenv('POSTGRES_HOST', 'localhost')),
                'port': db_secrets.get("POSTGRES_PORT", os.getenv('POSTGRES_PORT', '5432')),
                'database': db_secrets.get("POSTGRES_DATABASE", os.getenv('POSTGRES_DATABASE', 'postgres')),
                'user': db_secrets.get("POSTGRES_USER", os.getenv('POSTGRES_USER', 'postgres')),
                'password': db_secrets.get("POSTGRES_PASSWORD", os.getenv('POSTGRES_PASSWORD', ''))
            }
        else:
            raise FileNotFoundError("No secrets found")
    except (FileNotFoundError, KeyError, Exception):
        # Fall back to environment variables if no secrets file exists
        postgres_config = {
            'host': os.getenv('POSTGRES_HOST', 'localhost'),
            'port': os.getenv('POSTGRES_PORT', '5432'),
            'database': os.getenv('POSTGRES_DATABASE', 'postgres'),
            'user': os.getenv('POSTGRES_USER', 'postgres'),
            'password': os.getenv('POSTGRES_PASSWORD', '')
        }
    
    # Validate configuration
    if not all([postgres_config['host'], postgres_config['user'], postgres_config['password']]):
        st.error("⚠️ **Database configuration missing!** Please configure database credentials.")
        st.info("""
        **Local Development**: Create a .env file with your database credentials
        
        **Streamlit Community**: Add secrets in your app settings:
        ```
        [database]
        POSTGRES_HOST = "your-host"
        POSTGRES_PORT = "5432"
        POSTGRES_DATABASE = "your-database"
        POSTGRES_USER = "your-user"  
        POSTGRES_PASSWORD = "your-password"
        ```
        """)
        st.stop()
    
    # Initialize session state for activity tracking
    if 'last_activity' not in st.session_state:
        st.session_state.last_activity = datetime.now()
    if 'session_id' not in st.session_state:
        import uuid
        st.session_state.session_id = str(uuid.uuid4())

    # Update activity timestamp on each interaction
    st.session_state.last_activity = datetime.now()

    # Initialize extractor
    if 'extractor' not in st.session_state:
        with st.spinner("Initializing extraction system..."):
            st.session_state.extractor = OrganizationExtractor(postgres_config)
    
    # Check system status
    extractor = st.session_state.extractor
    if len(extractor.org_lookup) == 0:
        st.error("⚠️ **CRITICAL**: Lookup dictionary is empty! Database loading failed.")
        with st.spinner("Reloading database..."):
            extractor._load_organizations()
            extractor._build_lookup()
    
    extractor = st.session_state.extractor
    
    # Sidebar configuration
    st.sidebar.header("Settings")
    confidence_threshold = st.sidebar.slider(
        "Confidence Threshold",
        min_value=0.5,
        max_value=1.0,
        value=0.85,
        step=0.05
    )
    extractor.min_confidence = confidence_threshold

    # Keep-alive status in sidebar
    st.sidebar.header("Session Status")
    time_since_activity = datetime.now() - st.session_state.last_activity
    if time_since_activity.total_seconds() < 300:  # 5 minutes
        st.sidebar.success("🟢 Active")
    elif time_since_activity.total_seconds() < 900:  # 15 minutes
        st.sidebar.warning("🟡 Idle")
    else:
        st.sidebar.error("🔴 Long Idle")

    st.sidebar.caption(f"Last activity: {st.session_state.last_activity.strftime('%H:%M:%S')}")

    # Database statistics
    st.sidebar.metric("Organizations in Database", len(extractor.master_orgs_df))
    st.sidebar.metric("Lookup Entries", len(extractor.org_lookup))
    
    # File upload
    uploaded_files = st.file_uploader(
        "Upload Documents",
        type=['txt', 'pdf', 'docx', 'xlsx', 'xls', 'pptx'],
        help="Supported formats: TXT, PDF, DOCX, XLSX, XLS, PPTX",
        accept_multiple_files=True
    )
    
    # Clear cache button
    if st.button("🗑️ Clear Cache", help="Clear cached results to force re-processing"):
        # Clear extraction cache
        keys_to_remove = [key for key in st.session_state.keys() if key.startswith('extraction_')]
        for key in keys_to_remove:
            del st.session_state[key]
        # Clear approval states
        st.session_state.approved_orgs = []
        st.session_state.rejected_orgs = []
        st.session_state.manual_additions = []
        st.rerun()
    
    if uploaded_files:
        # Process all uploaded files
        all_file_results = {}
        all_db_matches = []
        all_ner_discoveries = []
        total_processing_time = 0
        total_characters = 0
        
        # Extract text and process each file
        for uploaded_file in uploaded_files:
            with st.spinner(f"Processing {uploaded_file.name}..."):
                # Extract text
                text = extract_text_from_file(uploaded_file)
                
                if text:
                    total_characters += len(text)
                    
                    # Extract organizations (cache results to avoid re-extraction)
                    # Use file content hash for more robust caching
                    import hashlib
                    text_hash = hashlib.md5(text.encode()).hexdigest()[:8]
                    cache_key = f"extraction_{uploaded_file.name}_{len(text)}_{text_hash}"
                    
                    if cache_key not in st.session_state:
                        start_time = time.time()
                        db_matches, ner_discoveries = extractor.extract_organizations(text)
                        processing_time = time.time() - start_time
                        
                        # Cache the results
                        st.session_state[cache_key] = {
                            'db_matches': db_matches,
                            'ner_discoveries': ner_discoveries,
                            'processing_time': processing_time
                        }
                    else:
                        # Use cached results
                        cached_results = st.session_state[cache_key]
                        db_matches = cached_results['db_matches']
                        ner_discoveries = cached_results['ner_discoveries'] 
                        processing_time = cached_results['processing_time']
                    
                    # Store results for this file
                    all_file_results[uploaded_file.name] = {
                        'text': text,
                        'db_matches': db_matches,
                        'ner_discoveries': ner_discoveries,
                        'processing_time': processing_time,
                        'characters': len(text)
                    }
                    
                    # Add to combined results
                    all_db_matches.extend(db_matches)
                    all_ner_discoveries.extend(ner_discoveries)
                    total_processing_time += processing_time
                else:
                    st.error(f"Could not extract text from {uploaded_file.name}")
        
        if all_file_results:
            # Summary section
            st.success(f"Successfully processed {len(all_file_results)} files")
            
            # Overall metrics
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Files Processed", len(all_file_results))
            with col2:
                st.metric("Total Characters", f"{total_characters:,}")
            with col3:
                st.metric("Total Processing Time", f"{total_processing_time:.2f}s")
            with col4:
                st.metric("Avg Time per File", f"{total_processing_time/len(all_file_results):.2f}s")
            
            # Deduplicate combined results
            db_matches = extractor._deduplicate_matches(all_db_matches)
            ner_discoveries = extractor._deduplicate_matches(all_ner_discoveries)
            
            # Combined results summary
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Combined Database Matches", len(db_matches))
            with col2:
                st.metric("Combined New Discoveries", len(ner_discoveries))
            
            # Debug section - only show in local development  
            # Set DEBUG_MODE=false in Streamlit Community Cloud secrets to hide debug info
            debug_mode = os.getenv('DEBUG_MODE', 'true').lower() == 'true'
            
            # Try to check secrets, but don't fail if no secrets file exists
            try:
                if 'DEBUG_MODE' in st.secrets:
                    debug_mode = st.secrets.get('DEBUG_MODE', 'true').lower() == 'true'
            except (FileNotFoundError, Exception):
                # No secrets file or other error - use environment variable or default
                pass
            
            if debug_mode:
                with st.expander("🔍 Debug Info", expanded=False):
                    st.subheader("Database Lookup Statistics")
                    st.write(f"Total organizations in lookup: {len(extractor.org_lookup)}")
                    st.write(f"Known short organizations from DB: {len(extractor.known_short_orgs)}")
                    
                    # Show a sample of lookup keys
                    sample_keys = list(extractor.org_lookup.keys())[:10]
                    st.write("Sample lookup keys:", sample_keys)
                    
                    # Show short orgs from database
                    if len(extractor.known_short_orgs) > 0:
                        short_orgs_sample = list(extractor.known_short_orgs)[:20]
                        st.write(f"Sample short orgs from DB: {short_orgs_sample}")
                        if len(extractor.known_short_orgs) > 20:
                            st.write(f"... and {len(extractor.known_short_orgs) - 20} more")
                    
                    # Show text samples from each file
                    st.subheader("Text Extraction Samples")
                    for filename, file_data in all_file_results.items():
                        with st.container():
                            st.write(f"**{filename}**: {len(file_data['text']):,} characters")
                            # Show first 500 chars to check extraction quality
                            preview = file_data['text'][:500].replace('\n', ' ')
                            st.text(f"Preview: {preview}...")
                    
                    # Manual org search
                    st.subheader("Manual Organization Search")
                    search_org = st.text_input("Search for specific org in extracted text:", key="debug_search")
                    
                    # Show all current matches for troubleshooting
                    st.subheader("All Current Matches (Troubleshooting)")
                    st.write("**Database Matches Found:**")
                    for i, match in enumerate(db_matches):
                        st.markdown(f"**DB Match {i+1}: {match['canonical']}**")
                        with st.container():
                            col1, col2 = st.columns([1, 2])
                            with col1:
                                st.json(match)
                            with col2:
                                # Show which file(s) this came from
                                source_files = []
                                for filename, file_data in all_file_results.items():
                                    for file_match in file_data['db_matches']:
                                        if file_match['canonical'] == match['canonical']:
                                            source_files.append(filename)
                                            # Show the actual text that was matched
                                            st.write(f"**Matched text in {filename}:** '{file_match['text']}'")
                                            
                                            # Show surrounding context
                                            file_text = file_data['text']
                                            match_pos = file_text.lower().find(file_match['text'].lower())
                                            if match_pos >= 0:
                                                start = max(0, match_pos - 100)
                                                end = min(len(file_text), match_pos + len(file_match['text']) + 100)
                                                context = file_text[start:end]
                                                highlighted = context.replace(file_match['text'], f"**{file_match['text']}**")
                                                st.write(f"**Context:** ...{highlighted}...")
                                            else:
                                                st.error("❌ Could not find this text in the file!")
                                
                                st.write(f"**Source files:** {list(set(source_files))}")
                        st.markdown("---")
                    
                    if ner_discoveries:
                        st.write("**NER Discoveries Found:**")
                        for i, discovery in enumerate(ner_discoveries):
                            st.markdown(f"**NER Discovery {i+1}: {discovery['text']}**")
                            
                            # Show fuzzy match debugging info
                            if 'fuzzy_score' in discovery:
                                if discovery['fuzzy_score'] > 0:
                                    st.write(f"🔍 **Fuzzy Match Debug**: Best match was '{discovery.get('best_fuzzy_match', 'Unknown')}' with score {discovery['fuzzy_score']}% (threshold: 85%)")
                                else:
                                    st.write(f"🔍 **Fuzzy Match Debug**: No fuzzy matches found")
                            
                            st.json(discovery)
                            st.markdown("---")
                    if search_org:
                        for filename, file_data in all_file_results.items():
                            text_lower = file_data['text'].lower()
                            search_lower = search_org.lower()
                            
                            # Check if it appears in text
                            if search_lower in text_lower:
                                st.success(f"✅ '{search_org}' found in {filename}")
                                
                                # Check if it's in lookup
                                if search_lower in extractor.org_lookup:
                                    st.info(f"✅ '{search_org}' exists in database lookup")
                                    
                                    # Test the actual regex matching logic
                                    import re
                                    pattern = r'\b' + re.escape(search_lower) + r'\b'
                                    matches = list(re.finditer(pattern, text_lower))
                                    
                                    if matches:
                                        st.success(f"✅ Regex pattern matches {len(matches)} time(s)")
                                        
                                        # Check if it would pass other filters
                                        org_info = extractor.org_lookup[search_lower]
                                        st.write(f"Organization info: {org_info}")
                                        
                                        # Check length filter
                                        if len(search_lower) >= extractor.min_org_length:
                                            st.success(f"✅ Passes length filter (≥{extractor.min_org_length})")
                                        else:
                                            st.error(f"❌ FAILS length filter (<{extractor.min_org_length})")
                                        
                                        # Check generic term filter
                                        if extractor._is_generic_term(search_lower):
                                            st.error(f"❌ FAILS generic term filter")
                                        else:
                                            st.success(f"✅ Passes generic term filter")
                                        
                                        # Check short term rules
                                        if len(search_lower) < 8 and ' ' not in search_lower:
                                            if search_lower not in extractor.known_short_orgs:
                                                st.error(f"❌ FAILS short term filter (not in known_short_orgs)")
                                            else:
                                                st.success(f"✅ Passes short term filter (in known_short_orgs)")
                                        else:
                                            st.success(f"✅ Passes short term filter (long enough or multi-word)")

                                        # Check false positive filtering
                                        if extractor._is_false_positive_prone(search_lower):
                                            st.warning(f"⚠️ Term is flagged as false-positive prone")
                                            penalty = extractor._get_confidence_penalty(search_lower)
                                            if penalty > 0:
                                                st.write(f"🎯 Confidence penalty: {penalty:.2%}")
                                            if extractor._requires_strong_context(search_lower):
                                                st.write(f"📋 Requires strong organizational context")
                                        else:
                                            st.success(f"✅ Not flagged as false-positive prone")
                                            
                                    else:
                                        st.error(f"❌ Regex pattern '{pattern}' does NOT match")
                                        # Show context around the term
                                        import re
                                        simple_matches = [m.start() for m in re.finditer(re.escape(search_lower), text_lower)]
                                        if simple_matches:
                                            st.write("Found at positions:", simple_matches[:5])
                                            for pos in simple_matches[:3]:
                                                start = max(0, pos - 50)
                                                end = min(len(text_lower), pos + len(search_lower) + 50)
                                                context = text_lower[start:end]
                                                highlighted = context.replace(search_lower, f"**{search_lower}**")
                                                st.write(f"Context: ...{highlighted}...")
                                    
                                else:
                                    st.warning(f"⚠️ '{search_org}' NOT in database lookup")
                                    
                                    # Try fuzzy search in lookup
                                    from fuzzywuzzy import process
                                    if len(extractor.org_lookup) > 0:
                                        matches = process.extract(search_org, extractor.org_lookup.keys(), limit=3)
                                        st.write(f"Closest matches in lookup: {matches}")
                            else:
                                st.error(f"❌ '{search_org}' NOT found in {filename}")
            
            # File-by-file breakdown
            with st.expander("📁 File-by-File Breakdown", expanded=False):
                for filename, file_data in all_file_results.items():
                    st.subheader(f"📄 {filename}")
                    
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("Characters", f"{file_data['characters']:,}")
                    with col2:
                        st.metric("DB Matches", len(file_data['db_matches']))
                    with col3:
                        st.metric("New Discoveries", len(file_data['ner_discoveries']))
                    with col4:
                        st.metric("Processing Time", f"{file_data['processing_time']:.2f}s")
                    
                    # File content preview
                    preview_text = file_data['text'][:500] + "..." if len(file_data['text']) > 500 else file_data['text']
                    st.text_area("Content Preview", preview_text, height=100, key=f"preview_{filename}")
                    
                    st.markdown("---")
            
            # Display results
            if db_matches or ner_discoveries:
                
                # Database Matches Tab
                st.subheader("✅ Organizations Found in Database")
                if db_matches:
                    db_df = pd.DataFrame(db_matches)
                    db_df['confidence'] = db_df['confidence'].round(3)
                    
                    # Allow filtering
                    selected_db = st.multiselect(
                        "Select confirmed organizations:",
                        options=range(len(db_df)),
                        default=list(range(len(db_df))),
                        format_func=lambda x: f"{db_df.iloc[x]['canonical']} ({db_df.iloc[x]['confidence']:.2f})"
                    )
                    
                    if selected_db:
                        confirmed_orgs = db_df.iloc[selected_db]
                        
                        # Show fuzzy matches separately for debugging
                        fuzzy_matches = confirmed_orgs[confirmed_orgs['method'] == 'ner_fuzzy']
                        if len(fuzzy_matches) > 0:
                            st.info(f"🔍 **{len(fuzzy_matches)} organizations were matched via fuzzy matching** (these were initially detected by NER but matched to existing database entries)")
                            
                            # Show display columns, including fuzzy_score if available
                            display_cols = ['canonical', 'confidence', 'method']
                            if 'fuzzy_score' in confirmed_orgs.columns:
                                display_cols.append('fuzzy_score')
                            
                            st.dataframe(confirmed_orgs[display_cols], use_container_width=True)
                        else:
                            st.dataframe(confirmed_orgs[['canonical', 'confidence', 'method']], use_container_width=True)
                else:
                    st.info("No organizations found in database")
                
                # New Discoveries Tab
                st.subheader("🔍 Potential New Organizations (Require Review)")
                if ner_discoveries:
                    st.warning(f"Found {len(ner_discoveries)} potential new organizations that need manual review")
                    
                    # Initialize session state for approved/rejected organizations
                    if 'approved_orgs' not in st.session_state:
                        st.session_state.approved_orgs = []
                    if 'rejected_orgs' not in st.session_state:
                        st.session_state.rejected_orgs = []
                    
                    for i, discovery in enumerate(ner_discoveries):
                        org_text = discovery['text']
                        
                        # Skip if already approved or rejected
                        if org_text in st.session_state.approved_orgs or org_text in st.session_state.rejected_orgs:
                            continue
                            
                        with st.container():
                            col1, col2, col3 = st.columns([3, 1, 1])
                            
                            with col1:
                                st.write(f"**{org_text}** (confidence: {discovery['confidence']:.2f})")
                                # Show fuzzy match info if available
                                if 'fuzzy_score' in discovery and discovery['fuzzy_score'] > 0:
                                    st.caption(f"🔍 Closest DB match: '{discovery.get('best_fuzzy_match', 'Unknown')}' ({discovery['fuzzy_score']}%)")
                            
                            with col2:
                                if st.button(f"Approve", key=f"approve_{i}"):
                                    st.session_state.approved_orgs.append(org_text)
                                    st.rerun()
                            
                            with col3:
                                if st.button(f"Reject", key=f"reject_{i}"):
                                    st.session_state.rejected_orgs.append(org_text)
                                    st.rerun()
                    
                    # Display approved organizations
                    if st.session_state.approved_orgs:
                        st.success(f"✅ **Approved Organizations ({len(st.session_state.approved_orgs)}):**")
                        for approved_org in st.session_state.approved_orgs:
                            st.write(f"• {approved_org}")
                    
                    # Display rejected organizations  
                    if st.session_state.rejected_orgs:
                        st.error(f"❌ **Rejected Organizations ({len(st.session_state.rejected_orgs)}):**")
                        for rejected_org in st.session_state.rejected_orgs:
                            st.write(f"• {rejected_org}")
                            
                    # Clear buttons
                    if st.session_state.approved_orgs or st.session_state.rejected_orgs:
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("🔄 Reset Approvals"):
                                st.session_state.approved_orgs = []
                                st.session_state.rejected_orgs = []
                                st.rerun()
                        with col2:
                            if st.button("📋 Add to Extracted List"):
                                # Add approved orgs to db_matches for export
                                added_count = 0
                                for approved_org in st.session_state.approved_orgs:
                                    db_matches.append({
                                        'text': approved_org,
                                        'canonical': approved_org,
                                        'confidence': 0.95,
                                        'org_id': None,
                                        'method': 'manual_approval'
                                    })
                                    added_count += 1
                                
                                # Add manual additions to db_matches for export
                                if 'manual_additions' in st.session_state:
                                    for manual_org in st.session_state.manual_additions:
                                        db_matches.append({
                                            'text': manual_org,
                                            'canonical': manual_org,
                                            'confidence': 0.90,
                                            'org_id': None,
                                            'method': 'manual_addition'
                                        })
                                        added_count += 1
                                
                                st.success(f"Added {added_count} organizations to extracted list!")
                                st.session_state.approved_orgs = []
                                if 'manual_additions' in st.session_state:
                                    st.session_state.manual_additions = []
                                st.rerun()
                else:
                    st.info("No new organizations discovered")
                
                # Manual Addition Section
                st.subheader("➕ Manual Organization Addition")
                
                # Initialize manual additions in session state
                if 'manual_additions' not in st.session_state:
                    st.session_state.manual_additions = []
                
                with st.form("manual_addition"):
                    new_org_name = st.text_input(
                        "Organization Name",
                        placeholder="Enter organization name not found by the system"
                    )
                    
                    submitted = st.form_submit_button("Add to Pending List")
                    
                    if submitted and new_org_name:
                        new_org_name = new_org_name.strip()
                        if new_org_name.lower() in extractor.org_lookup:
                            st.warning(f"'{new_org_name}' already exists in database")
                        elif new_org_name in st.session_state.manual_additions:
                            st.warning(f"'{new_org_name}' already added to pending list")
                        else:
                            st.session_state.manual_additions.append(new_org_name)
                            st.success(f"Added '{new_org_name}' to pending additions")
                            st.rerun()
                
                # Display current manual additions
                if st.session_state.manual_additions:
                    st.success(f"📋 **Manually Added Organizations ({len(st.session_state.manual_additions)}):**")
                    for i, manual_org in enumerate(st.session_state.manual_additions):
                        col1, col2 = st.columns([4, 1])
                        with col1:
                            st.write(f"• {manual_org}")
                        with col2:
                            if st.button("❌", key=f"remove_manual_{i}", help="Remove from list"):
                                st.session_state.manual_additions.remove(manual_org)
                                st.rerun()
                    
                    if st.button("🔄 Clear All Manual Additions"):
                        st.session_state.manual_additions = []
                        st.rerun()
                
                # Export Results
                st.subheader("📥 Export Combined Results")
                
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    if db_matches:
                        # Add source file information to CSV export
                        enriched_db_matches = []
                        for match in db_matches:
                            # Find which file(s) this match came from
                            source_files = []
                            for filename, file_data in all_file_results.items():
                                for file_match in file_data['db_matches']:
                                    if file_match['canonical'] == match['canonical']:
                                        source_files.append(filename)
                            
                            match_with_source = match.copy()
                            match_with_source['source_files'] = '; '.join(set(source_files))
                            enriched_db_matches.append(match_with_source)
                        
                        csv_data = pd.DataFrame(enriched_db_matches).to_csv(index=False)
                        st.download_button(
                            label="📊 CSV (DB Matches)",
                            data=csv_data,
                            file_name=f"db_matches_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                            mime="text/csv"
                        )
                
                with col2:
                    if ner_discoveries:
                        # Add source file information to NER discoveries
                        enriched_ner_discoveries = []
                        for discovery in ner_discoveries:
                            # Find which file(s) this discovery came from
                            source_files = []
                            for filename, file_data in all_file_results.items():
                                for file_discovery in file_data['ner_discoveries']:
                                    if file_discovery['text'] == discovery['text']:
                                        source_files.append(filename)
                            
                            discovery_with_source = discovery.copy()
                            discovery_with_source['source_files'] = '; '.join(set(source_files))
                            enriched_ner_discoveries.append(discovery_with_source)
                        
                        csv_data = pd.DataFrame(enriched_ner_discoveries).to_csv(index=False)
                        st.download_button(
                            label="🔍 CSV (New Orgs)",
                            data=csv_data,
                            file_name=f"new_orgs_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                            mime="text/csv"
                        )
                
                with col3:
                    if db_matches or (st.session_state.get('approved_orgs', [])) or (st.session_state.get('manual_additions', [])):
                        # Create combined organization list for .txt export (DB matches + approved orgs + manual additions)
                        all_orgs = []
                        
                        # Add database matches
                        for match in db_matches:
                            all_orgs.append(match['canonical'])
                        
                        # Add approved organizations from review process
                        if 'approved_orgs' in st.session_state:
                            all_orgs.extend(st.session_state.approved_orgs)
                        
                        # Add manual additions
                        if 'manual_additions' in st.session_state:
                            all_orgs.extend(st.session_state.manual_additions)
                        
                        # Remove duplicates and sort
                        unique_orgs = sorted(list(set(all_orgs)))
                        
                        # Show what will be exported
                        db_count = len([m['canonical'] for m in db_matches])
                        approved_count = len(st.session_state.get('approved_orgs', []))
                        manual_count = len(st.session_state.get('manual_additions', []))
                        st.write(f"**Export will include:** {db_count} DB matches + {approved_count} approved orgs + {manual_count} manual additions = {len(unique_orgs)} total")
                        
                        # Separator selection
                        separator = st.radio(
                            "Choose separator:",
                            options=[", ", "; "],
                            format_func=lambda x: "Comma" if x == ", " else "Semicolon",
                            horizontal=True,
                            key="separator_choice"
                        )
                        
                        # Create separated string
                        txt_content = separator.join(unique_orgs)
                        
                        st.download_button(
                            label="📝 TXT (Approved Orgs)",
                            data=txt_content,
                            file_name=f"approved_organizations_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                            mime="text/plain"
                        )
                
                with col4:
                    # Create comprehensive report with file breakdown
                    if all_file_results:
                        report_lines = []
                        report_lines.append(f"ORGANIZATION EXTRACTION REPORT")
                        report_lines.append(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                        report_lines.append(f"Files Processed: {len(all_file_results)}")
                        report_lines.append(f"Total Characters: {total_characters:,}")
                        report_lines.append(f"Combined DB Matches: {len(db_matches)}")
                        report_lines.append(f"Combined New Discoveries: {len(ner_discoveries)}")
                        report_lines.append("")
                        
                        # File breakdown
                        report_lines.append("FILE BREAKDOWN:")
                        for filename, file_data in all_file_results.items():
                            report_lines.append(f"\n{filename}:")
                            report_lines.append(f"  Characters: {file_data['characters']:,}")
                            report_lines.append(f"  DB Matches: {len(file_data['db_matches'])}")
                            report_lines.append(f"  New Discoveries: {len(file_data['ner_discoveries'])}")
                            report_lines.append(f"  Processing Time: {file_data['processing_time']:.2f}s")
                        
                        # Combined organizations
                        report_lines.append("\n\nCOMBINED ORGANIZATIONS:")
                        all_orgs = []
                        for match in db_matches:
                            all_orgs.append(match['canonical'])
                        for discovery in ner_discoveries:
                            all_orgs.append(discovery['text'])
                        if 'manual_additions' in st.session_state:
                            all_orgs.extend(st.session_state.manual_additions)
                        unique_orgs = sorted(list(set(all_orgs)))
                        report_lines.extend([f"- {org}" for org in unique_orgs])
                        
                        # Manual additions section
                        if 'manual_additions' in st.session_state and st.session_state.manual_additions:
                            report_lines.append(f"\n\nMANUAL ADDITIONS ({len(st.session_state.manual_additions)}):")
                            for manual_org in st.session_state.manual_additions:
                                report_lines.append(f"- {manual_org}")
                        
                        report_content = "\n".join(report_lines)
                        
                        st.download_button(
                            label="📄 Full Report",
                            data=report_content,
                            file_name=f"org_extraction_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                            mime="text/plain"
                        )
            
            else:
                st.warning("No organizations found in the document")
    
    # Footer
    st.markdown("---")
    st.markdown("**Note:** This tool uses BERT-based NER combined with database matching for organization extraction.")


if __name__ == "__main__":
    main()